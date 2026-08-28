from __future__ import annotations

import argparse
import hashlib
import html
import json
from pathlib import Path

from docx import Document
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_FIXTURE = ROOT / "tests" / "fixtures" / "youtube-learning" / "pcr-qpcr-course-large-golden.json"
DEFAULT_OUT = ROOT / "build" / "youtube-course-large-golden"


def canonical_fingerprint(data: dict) -> str:
    payload = json.dumps(data, ensure_ascii=False, sort_keys=True, separators=(",", ":")).encode("utf-8")
    return hashlib.sha256(payload).hexdigest()


def course_model(data: dict, fingerprint: str) -> dict:
    return {
        "schemaVersion": 1,
        "courseId": data["courseId"],
        "title": data["title"],
        "audience": data["audience"],
        "language": data["language"],
        "courseDepth": data["courseDepth"],
        "courseFingerprint": fingerprint,
        "sources": data["sources"],
        "assumedPrerequisites": data["assumedPrerequisites"],
        "modules": data["modules"],
        "learningObjectives": data["learningObjectives"],
        "knowledgeChecks": data["knowledgeChecks"],
        "designAuthority": data["designAuthority"],
        "renderTargets": data["renderTargets"],
    }


def generate_svg(data: dict, fingerprint: str, out: Path) -> None:
    modules = data["modules"]
    w, h = 1400, 250 + len(modules) * 125
    boxes = []
    edges = []
    y_by_id = {}
    for i, module in enumerate(modules):
        y = 150 + i * 125
        y_by_id[module["moduleId"]] = y
        title = html.escape(module["title"])
        promise = html.escape(module["competencePromise"])
        boxes.append(
            f'<rect x="270" y="{y}" width="860" height="88" rx="14" fill="white" stroke="#343a40" stroke-width="2"/>'
            f'<text x="300" y="{y+31}" font-family="Arial, sans-serif" font-size="22" font-weight="700">{module["moduleId"]} · {title}</text>'
            f'<text x="300" y="{y+62}" font-family="Arial, sans-serif" font-size="16">{promise}</text>'
        )
    for module in modules:
        for prereq in module.get("prerequisites", []):
            y1 = y_by_id[prereq] + 88
            y2 = y_by_id[module["moduleId"]]
            edges.append(f'<path d="M700 {y1} L700 {y2}" stroke="#343a40" stroke-width="2.5" marker-end="url(#arrow)"/>')
    svg = f'''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {w} {h}" role="img" aria-labelledby="title desc">
<title id="title">Course Map PCR und qPCR</title>
<desc id="desc">Acht Module mit fachlichen Voraussetzungskanten. Course fingerprint {fingerprint[:12]}.</desc>
<defs><marker id="arrow" markerWidth="10" markerHeight="10" refX="8" refY="3" orient="auto"><path d="M0,0 L0,6 L9,3 z" fill="#343a40"/></marker></defs>
<rect width="100%" height="100%" fill="#f8f9fa"/>
<text x="90" y="70" font-family="Arial, sans-serif" font-size="34" font-weight="700">PCR → qPCR Learning Path</text>
<text x="90" y="105" font-family="Arial, sans-serif" font-size="16">Course fingerprint: {fingerprint[:12]} · required path based on prerequisites, not playlist order</text>
{''.join(edges)}
{''.join(boxes)}
</svg>'''
    out.write_text(svg, encoding="utf-8")


def generate_html(data: dict, fingerprint: str, out: Path) -> None:
    module_sections = []
    for module in data["modules"]:
        objectives = "".join(f"<li>{html.escape(data['learningObjectives'][oid])}</li>" for oid in module["objectives"])
        sources = ", ".join(module["sourceScope"])
        module_sections.append(f'''<section class="module" id="{module['moduleId']}">
<h2>{module['moduleId']} · {html.escape(module['title'])}</h2>
<p class="promise">{html.escape(module['competencePromise'])}</p>
<h3>Lernziele</h3><ul>{objectives}</ul>
<p><strong>Voraussetzungen:</strong> {html.escape(', '.join(module.get('prerequisites', [])) or 'Einstieg')}</p>
<p><strong>Quellen:</strong> {html.escape(sources)}</p>
</section>''')
    checks = "".join(f"<li><strong>{q['questionId']}</strong> {html.escape(q['prompt'])}</li>" for q in data["knowledgeChecks"])
    source_rows = "".join(
        f'<tr><td>{s["sourceId"]}</td><td>{html.escape(s["channel"])}</td><td><a href="https://www.youtube.com/watch?v={s["youtubeVideoId"]}">{html.escape(s["title"])}</a></td><td>{html.escape(s["role"])}</td></tr>'
        for s in data["sources"]
    )
    content = f'''<!doctype html><html lang="de"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>{html.escape(data['title'])}</title>
<style>body{{font-family:Arial,sans-serif;margin:0;color:#212529;background:#f8f9fa}}main{{max-width:1120px;margin:auto;padding:32px}}header{{padding:52px 32px;background:#fff;border-bottom:1px solid #dee2e6}}header>div{{max-width:1120px;margin:auto}}h1{{font-size:42px;margin:0 0 12px}}.meta{{color:#495057}}.course-map{{width:100%;background:white;margin:28px 0;border:1px solid #dee2e6}}.grid{{display:grid;grid-template-columns:repeat(2,minmax(0,1fr));gap:20px}}.module{{background:white;border:1px solid #dee2e6;border-radius:14px;padding:24px}}.promise{{font-size:18px}}table{{width:100%;border-collapse:collapse;background:white}}th,td{{padding:10px;border:1px solid #dee2e6;text-align:left}}@media(max-width:760px){{h1{{font-size:32px}}.grid{{grid-template-columns:1fr}}main{{padding:18px}}}}@media print{{body{{background:white}}.module{{break-inside:avoid}}}}</style></head>
<body data-course-fingerprint="{fingerprint}"><header><div><p>Golden Course Reference · 12 reale YouTube-Quellen</p><h1>{html.escape(data['title'])}</h1><p class="meta">Course fingerprint {fingerprint[:12]} · Zielgruppe: {html.escape(data['audience'])} · Tiefe: {html.escape(data['courseDepth'])}</p></div></header>
<main><h2>Course Map</h2><object class="course-map" data="course-map.svg" type="image/svg+xml" aria-label="Course Map"></object><div class="grid">{''.join(module_sections)}</div>
<section><h2>Formative Knowledge Checks</h2><ol>{checks}</ol></section><section><h2>Source Map</h2><table><thead><tr><th>ID</th><th>Quelle</th><th>Video</th><th>Rolle</th></tr></thead><tbody>{source_rows}</tbody></table></section></main></body></html>'''
    out.write_text(content, encoding="utf-8")


def generate_pptx(data: dict, fingerprint: str, out: Path) -> None:
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(PInches(0.8), PInches(1.0), PInches(11.7), PInches(4.5)).text_frame
    p = tb.paragraphs[0]; p.text = data["title"]; p.font.size = PPt(30); p.font.bold = True
    p = tb.add_paragraph(); p.text = "Instructor / Workshop Golden Reference"; p.font.size = PPt(20)
    p = tb.add_paragraph(); p.text = f"12 sources · 8 modules · fingerprint {fingerprint[:12]}"; p.font.size = PPt(15)
    for module in data["modules"]:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tf = slide.shapes.add_textbox(PInches(0.7), PInches(0.55), PInches(12.0), PInches(6.1)).text_frame
        p = tf.paragraphs[0]; p.text = f"{module['moduleId']} · {module['title']}"; p.font.size = PPt(28); p.font.bold = True
        p = tf.add_paragraph(); p.text = module["competencePromise"]; p.font.size = PPt(18)
        p = tf.add_paragraph(); p.text = "Lernziele"; p.font.size = PPt(18); p.font.bold = True
        for oid in module["objectives"]:
            p = tf.add_paragraph(); p.text = data["learningObjectives"][oid]; p.level = 1; p.font.size = PPt(16)
        p = tf.add_paragraph(); p.text = f"Quellen: {', '.join(module['sourceScope'])}"; p.font.size = PPt(13)
        p = tf.add_paragraph(); p.text = f"Exit Criteria: {'; '.join(module['exitCriteria'])}"; p.font.size = PPt(13)
        footer = slide.shapes.add_textbox(PInches(10.6), PInches(7.05), PInches(2.1), PInches(0.25)).text_frame.paragraphs[0]
        footer.text = fingerprint[:12]; footer.font.size = PPt(8)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tf = slide.shapes.add_textbox(PInches(0.8), PInches(0.6), PInches(11.7), PInches(6.1)).text_frame
    p = tf.paragraphs[0]; p.text = "Formative Knowledge Checks"; p.font.size = PPt(28); p.font.bold = True
    for q in data["knowledgeChecks"]:
        p = tf.add_paragraph(); p.text = f"{q['questionId']}: {q['prompt']}"; p.font.size = PPt(14)
    prs.core_properties.subject = f"Course fingerprint {fingerprint}"
    prs.save(out)


def generate_docx(data: dict, fingerprint: str, out: Path) -> None:
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.7); sec.bottom_margin = Inches(0.7); sec.left_margin = Inches(0.8); sec.right_margin = Inches(0.8)
    doc.add_heading(data["title"], 0)
    doc.add_paragraph(f"Study Guide · 12 sources · 8 modules · Course fingerprint {fingerprint[:12]}")
    doc.add_heading("Course Map", level=1)
    for module in data["modules"]:
        prereq = ", ".join(module.get("prerequisites", [])) or "Einstieg"
        doc.add_paragraph(f"{module['moduleId']} → {module['title']} | Voraussetzung: {prereq}")
    for module in data["modules"]:
        doc.add_page_break()
        doc.add_heading(f"{module['moduleId']} · {module['title']}", level=1)
        doc.add_paragraph(module["competencePromise"])
        doc.add_heading("Lernziele", level=2)
        for oid in module["objectives"]:
            doc.add_paragraph(data["learningObjectives"][oid], style="List Bullet")
        doc.add_heading("Exit Criteria", level=2)
        for item in module["exitCriteria"]:
            doc.add_paragraph(item, style="List Bullet")
        doc.add_paragraph(f"Quellen: {', '.join(module['sourceScope'])}")
        related = [q for q in data["knowledgeChecks"] if q["moduleId"] == module["moduleId"]]
        if related:
            doc.add_heading("Checkpoint", level=2)
            for q in related:
                doc.add_paragraph(q["prompt"])
    doc.add_page_break(); doc.add_heading("Source Map", level=1)
    table = doc.add_table(rows=1, cols=4)
    hdr = table.rows[0].cells
    hdr[0].text = "ID"; hdr[1].text = "Channel"; hdr[2].text = "Video"; hdr[3].text = "Role"
    for s in data["sources"]:
        cells = table.add_row().cells
        cells[0].text = s["sourceId"]; cells[1].text = s["channel"]; cells[2].text = s["title"]; cells[3].text = s["role"]
    footer = sec.footer.paragraphs[0]
    footer.text = f"Golden Reference · {fingerprint[:12]}"
    footer.style.font = None if not hasattr(footer.style, "font") else footer.style.font
    for run in footer.runs:
        run.font.size = Pt(8)
    doc.core_properties.subject = f"Course fingerprint {fingerprint}"
    doc.save(out)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--fixture", type=Path, default=DEFAULT_FIXTURE)
    parser.add_argument("--out", type=Path, default=DEFAULT_OUT)
    args = parser.parse_args()
    data = json.loads(args.fixture.read_text(encoding="utf-8"))
    fingerprint = canonical_fingerprint(data)
    args.out.mkdir(parents=True, exist_ok=True)
    model = course_model(data, fingerprint)
    (args.out / "course-learning-model.json").write_text(json.dumps(model, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    generate_svg(data, fingerprint, args.out / "course-map.svg")
    generate_html(data, fingerprint, args.out / "index.html")
    generate_pptx(data, fingerprint, args.out / "instructor-deck.pptx")
    generate_docx(data, fingerprint, args.out / "study-guide.docx")
    manifest = {"courseFingerprint": fingerprint, "sourceCount": len(data["sources"]), "moduleCount": len(data["modules"]), "artifacts": ["course-learning-model.json","course-map.svg","index.html","instructor-deck.pptx","study-guide.docx"]}
    (args.out / "artifact-manifest.json").write_text(json.dumps(manifest, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(manifest, indent=2))


if __name__ == "__main__":
    main()
